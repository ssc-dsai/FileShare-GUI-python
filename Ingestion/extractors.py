# Ingestion/extractors.py
"""
Text extraction + multi-image export for vision.
Pure Python – no OCR, no Ollama.

All vision assets (embedded + standalone) are stored under:
  EXTRACTED_TEXTS_DIR / _images / {stem} /
Folders are created only when at least one image is actually saved.
"""

from __future__ import annotations

import logging
import zipfile
from pathlib import Path

logger = logging.getLogger(__name__)

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".tif", ".tiff", ".bmp", ".webp"}


def _images_dir_for(path: Path, extracted_texts_dir: Path | None) -> Path:
    """Return target folder path. Does not create the directory."""
    stem = path.stem
    if extracted_texts_dir is not None:
        return Path(extracted_texts_dir) / "_images" / stem
    return path.parent / "_images_fallback" / stem


def _ensure_dir(d: Path) -> None:
    d.mkdir(parents=True, exist_ok=True)


def _vision_header(image_paths: list[Path]) -> str:
    if not image_paths:
        return "[VISION_FLAG: No]\n"
    lines = [
        "[VISION_FLAG: Yes]",
        f"[IMAGE_COUNT: {len(image_paths)}]",
    ]
    for p in image_paths:
        lines.append(f"[IMAGE_PATH: {p.resolve()}]")
    return "\n".join(lines) + "\n"


def extract_text_from_file(
    file_path: str | Path,
    extracted_texts_dir: str | Path | None = None,
) -> str:
    """
    Extract text and export embedded/standalone images for later vision.
    Always returns a string that begins with [VISION_FLAG: …].
    """
    path = Path(file_path)
    if not path.is_file():
        raise FileNotFoundError(f"File not found: {path}")

    ext = path.suffix.lower()
    et_dir = Path(extracted_texts_dir) if extracted_texts_dir else None

    try:
        if ext in IMAGE_EXTENSIONS:
            # Standalone image → copy into extracted_texts/_images/{stem}/
            image_paths: list[Path] = [path]
            if et_dir is not None:
                out_dir = _images_dir_for(path, et_dir)
                _ensure_dir(out_dir)
                dest = out_dir / path.name
                if not dest.exists() or dest.stat().st_size != path.stat().st_size:
                    dest.write_bytes(path.read_bytes())
                image_paths = [dest]
                logger.info(f"{path.name}: standalone image copied → {dest}")
            return (
                _vision_header(image_paths)
                + "\n[IMAGE – pixel content reserved for vision model; no OCR text]\n"
            )

        if ext == ".pdf":
            text, imgs = _extract_pdf(path, et_dir)
        elif ext in {".doc", ".docx"}:
            text, imgs = _extract_docx(path, et_dir)
        elif ext in {".ppt", ".pptx"}:
            text, imgs = _extract_pptx(path, et_dir)
        elif ext == ".txt":
            text = path.read_text(encoding="utf-8", errors="replace").strip()
            imgs = []
        else:
            raise ValueError(f"Unsupported extension: {ext}")

        return _vision_header(imgs) + "\n" + (text or "").strip()

    except Exception as e:
        logger.error(f"Extraction failed {path.name}: {type(e).__name__}: {e}")
        return "[VISION_FLAG: No]\n[EXTRACTION FAILED]\n"


def _save_pdf_pixmap(doc, xref: int, out_path: Path) -> bool:
    """Save one PDF embedded image. Returns True if a usable file was written."""
    if fitz is None:
        return False
    try:
        pix = fitz.Pixmap(doc, xref)
        if pix.n - pix.alpha > 3:
            pix = fitz.Pixmap(fitz.csRGB, pix)
        elif pix.alpha:
            pix = fitz.Pixmap(fitz.csRGB, pix)
        if pix.width < 20 or pix.height < 20:
            return False
        _ensure_dir(out_path.parent)
        pix.save(str(out_path))
        return out_path.is_file() and out_path.stat().st_size > 0
    except Exception as e:
        logger.warning(f"PDF image export failed xref={xref}: {e}")
        return False


def _extract_pdf(path: Path, extracted_texts_dir: Path | None) -> tuple[str, list[Path]]:
    if fitz is None:
        return "[PDF EXTRACTION FAILED – PyMuPDF missing]", []

    text_parts: list[str] = []
    image_paths: list[Path] = []
    out_dir = _images_dir_for(path, extracted_texts_dir)
    img_counter = 0

    doc = fitz.open(path)
    try:
        for page_i, page in enumerate(doc):
            page_text = page.get_text("text").strip()
            if page_text:
                text_parts.append(page_text)

            for img_meta in page.get_images(full=True):
                xref = img_meta[0]
                candidate = img_counter + 1
                out_path = out_dir / f"img_{candidate:03d}_p{page_i + 1}.png"
                if _save_pdf_pixmap(doc, xref, out_path):
                    img_counter = candidate
                    image_paths.append(out_path)
                    text_parts.append(
                        f"[IMAGE_REF: page={page_i + 1} index={img_counter} file={out_path.name}]"
                    )
    finally:
        doc.close()

    logger.info(f"{path.name}: PDF images exported = {len(image_paths)}")
    return "\n\n".join(text_parts), image_paths


def _extract_docx(path: Path, extracted_texts_dir: Path | None) -> tuple[str, list[Path]]:
    if Document is None:
        return "[DOCX EXTRACTION FAILED – python-docx missing]", []

    doc = Document(str(path))
    paras = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    text = "\n".join(paras)

    image_paths: list[Path] = []
    out_dir = _images_dir_for(path, extracted_texts_dir)

    try:
        with zipfile.ZipFile(path, "r") as zf:
            media = [n for n in zf.namelist() if n.startswith("word/media/")]
            idx = 0
            for name in sorted(media):
                suffix = Path(name).suffix.lower() or ".bin"
                if suffix not in IMAGE_EXTENSIONS:
                    continue
                idx += 1
                out_path = out_dir / f"img_{idx:03d}{suffix}"
                data = zf.read(name)
                if not data:
                    idx -= 1
                    continue
                _ensure_dir(out_dir)
                out_path.write_bytes(data)
                if out_path.is_file() and out_path.stat().st_size > 0:
                    image_paths.append(out_path)
                else:
                    idx -= 1
    except Exception as e:
        logger.warning(f"DOCX media extract failed {path.name}: {e}")

    logger.info(f"{path.name}: DOCX images exported = {len(image_paths)}")
    return text, image_paths


def _extract_pptx(path: Path, extracted_texts_dir: Path | None) -> tuple[str, list[Path]]:
    if Presentation is None:
        return "[PPTX EXTRACTION FAILED – python-pptx missing]", []

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    parts.append(t)
    text = "\n".join(parts)

    image_paths: list[Path] = []
    out_dir = _images_dir_for(path, extracted_texts_dir)

    try:
        with zipfile.ZipFile(path, "r") as zf:
            media = [n for n in zf.namelist() if n.startswith("ppt/media/")]
            idx = 0
            for name in sorted(media):
                suffix = Path(name).suffix.lower() or ".bin"
                if suffix not in IMAGE_EXTENSIONS:
                    continue
                idx += 1
                out_path = out_dir / f"img_{idx:03d}{suffix}"
                data = zf.read(name)
                if not data:
                    idx -= 1
                    continue
                _ensure_dir(out_dir)
                out_path.write_bytes(data)
                if out_path.is_file() and out_path.stat().st_size > 0:
                    image_paths.append(out_path)
                else:
                    idx -= 1
    except Exception as e:
        logger.warning(f"PPTX media extract failed {path.name}: {e}")

    logger.info(f"{path.name}: PPTX images exported = {len(image_paths)}")
    return text, image_paths